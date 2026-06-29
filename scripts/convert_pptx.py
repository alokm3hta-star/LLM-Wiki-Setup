"""Convert a PPTX file to Markdown. Writes result to stdout.

Each slide becomes:
  ## [Slide Title]
  [Body text paragraphs]
  > [Speaker notes]
"""
import sys


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 convert_pptx.py <path-to-pptx>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]

    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        print(
            "ERROR: python-pptx is not installed. Run: pip install python-pptx",
            file=sys.stderr,
        )
        sys.exit(2)

    prs = Presentation(path)
    output = []

    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_parts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:  # Picture
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Heuristic: first non-empty text shape with title placeholder is the title
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title_text = text
            else:
                body_parts.append(text)

        if not title_text:
            title_text = f"Slide {i}"

        output.append(f"## {title_text}")
        for part in body_parts:
            output.append(part)
            output.append("")

        # Speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                for note_line in notes.splitlines():
                    output.append(f"> {note_line}")
                output.append("")

        output.append("")

    print("\n".join(output))


if __name__ == "__main__":
    main()
